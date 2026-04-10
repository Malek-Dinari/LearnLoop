import pytest
from app.services.document_service import DocumentService


@pytest.fixture
def doc_service():
    return DocumentService()


def test_chunk_text_basic(doc_service):
    text = "Paragraph one.\n\nParagraph two.\n\nParagraph three."
    chunks = doc_service.chunk_text(text, chunk_size=50, overlap=0)
    assert len(chunks) >= 1
    for chunk in chunks:
        assert len(chunk) > 0


def test_chunk_text_single(doc_service):
    text = "Short text."
    chunks = doc_service.chunk_text(text, chunk_size=1000, overlap=0)
    assert len(chunks) == 1
    assert "Short text" in chunks[0]


def test_chunk_text_long_paragraph(doc_service):
    text = " ".join(["word"] * 500)
    chunks = doc_service.chunk_text(text, chunk_size=100, overlap=0)
    assert len(chunks) > 1


def test_get_document_summary(doc_service):
    chunks = ["chunk one", "chunk two", "chunk three"]
    summary = doc_service.get_document_summary(chunks)
    assert "chunk one" in summary


@pytest.mark.asyncio
async def test_process_upload_txt(doc_service, tmp_path):
    doc_service.upload_dir = str(tmp_path)
    result = await doc_service.process_upload("test.txt", b"Hello world. This is test content.")
    assert "document_id" in result
    assert result["filename"] == "test.txt"
    assert result["chunk_count"] >= 1


@pytest.mark.asyncio
async def test_process_upload_docx(doc_service, tmp_path):
    from docx import Document as DocxDocument

    doc_service.upload_dir = str(tmp_path)
    # Create a minimal DOCX in memory
    doc = DocxDocument()
    doc.add_paragraph("This is a test paragraph about machine learning.")
    doc.add_paragraph("Neural networks are a key concept.")
    docx_path = tmp_path / "test.docx"
    doc.save(str(docx_path))
    content = docx_path.read_bytes()

    result = await doc_service.process_upload("test.docx", content)
    assert "document_id" in result
    assert result["filename"] == "test.docx"
    assert result["chunk_count"] >= 1


@pytest.mark.asyncio
async def test_process_upload_pptx(doc_service, tmp_path):
    from pptx import Presentation

    doc_service.upload_dir = str(tmp_path)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Test Slide"
    slide.placeholders[1].text = "This is test content for PPTX extraction."
    pptx_path = tmp_path / "test.pptx"
    prs.save(str(pptx_path))
    content = pptx_path.read_bytes()

    result = await doc_service.process_upload("test.pptx", content)
    assert "document_id" in result
    assert result["filename"] == "test.pptx"
    assert result["chunk_count"] >= 1


@pytest.mark.asyncio
async def test_extract_corrupt_docx(doc_service, tmp_path):
    from app.exceptions import DocumentProcessingError

    doc_service.upload_dir = str(tmp_path)
    with pytest.raises(DocumentProcessingError):
        await doc_service.process_upload("bad.docx", b"not-a-real-docx")


@pytest.mark.asyncio
async def test_extract_empty_text(doc_service, tmp_path):
    from app.exceptions import DocumentProcessingError

    doc_service.upload_dir = str(tmp_path)
    with pytest.raises(DocumentProcessingError, match="no extractable text"):
        await doc_service.process_upload("empty.txt", b"   ")
