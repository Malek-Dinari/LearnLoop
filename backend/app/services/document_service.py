import logging
import os
import uuid
from typing import TYPE_CHECKING

import fitz  # PyMuPDF

from app.config import settings
from app.exceptions import DocumentProcessingError

if TYPE_CHECKING:
    from sqlalchemy.ext.asyncio import AsyncSession

logger = logging.getLogger(__name__)


class DocumentService:
    def __init__(self):
        self.upload_dir = settings.upload_dir
        # document_id -> {filename, file_path, chunks}
        self.documents: dict[str, dict] = {}

    async def save_file(self, filename: str, content: bytes) -> tuple[str, str]:
        os.makedirs(self.upload_dir, exist_ok=True)
        document_id = str(uuid.uuid4())
        ext = os.path.splitext(filename)[1]
        file_path = os.path.join(self.upload_dir, f"{document_id}{ext}")
        with open(file_path, "wb") as f:
            f.write(content)
        return document_id, file_path

    async def extract_text(self, file_path: str, file_type: str) -> str:
        try:
            if file_type == "pdf":
                return self._extract_pdf(file_path)
            elif file_type == "docx":
                return self._extract_docx(file_path)
            elif file_type == "pptx":
                return self._extract_pptx(file_path)
            else:
                return self._extract_txt(file_path)
        except DocumentProcessingError:
            raise
        except Exception as exc:
            raise DocumentProcessingError(
                f"Failed to extract text from {file_type.upper()} file: {exc}"
            ) from exc

    def _extract_pdf(self, file_path: str) -> str:
        try:
            doc = fitz.open(file_path)
        except Exception as exc:
            raise DocumentProcessingError(
                "Could not read PDF — file may be corrupt or password-protected"
            ) from exc
        try:
            text = "\n\n".join([page.get_text() for page in doc])
        finally:
            doc.close()
        return text

    def _extract_txt(self, file_path: str) -> str:
        # Try UTF-8 first, then latin-1 fallback (covers most Western encodings)
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except UnicodeDecodeError:
            logger.info("UTF-8 decode failed for %s, trying latin-1", file_path)
            try:
                with open(file_path, "r", encoding="latin-1") as f:
                    return f.read()
            except Exception as exc:
                raise DocumentProcessingError(
                    "Could not decode text file — unsupported encoding"
                ) from exc

    def _extract_docx(self, file_path: str) -> str:
        from docx import Document as DocxDocument
        doc = DocxDocument(file_path)
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())

    def _extract_pptx(self, file_path: str) -> str:
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text.strip())
        return "\n\n".join(texts)

    def chunk_text(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> list[str]:
        paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]

        if not paragraphs:
            paragraphs = [p.strip() for p in text.split("\n") if p.strip()]

        if not paragraphs:
            return [text[:chunk_size]] if text.strip() else []

        chunks = []
        current_chunk = ""

        for para in paragraphs:
            if len(current_chunk) + len(para) + 2 <= chunk_size:
                current_chunk = f"{current_chunk}\n\n{para}" if current_chunk else para
            else:
                if current_chunk:
                    chunks.append(current_chunk)
                # If paragraph itself is too large, split it
                if len(para) > chunk_size:
                    words = para.split()
                    current_chunk = ""
                    for word in words:
                        if len(current_chunk) + len(word) + 1 <= chunk_size:
                            current_chunk = f"{current_chunk} {word}" if current_chunk else word
                        else:
                            if current_chunk:
                                chunks.append(current_chunk)
                            current_chunk = word
                else:
                    current_chunk = para

        if current_chunk:
            chunks.append(current_chunk)

        # Add overlap between chunks
        if len(chunks) > 1 and overlap > 0:
            overlapped = [chunks[0]]
            for i in range(1, len(chunks)):
                prev_tail = chunks[i - 1][-overlap:]
                overlapped.append(prev_tail + "\n\n" + chunks[i])
            chunks = overlapped

        return chunks

    def get_document_summary(self, chunks: list[str]) -> str:
        combined = "\n\n".join(chunks)
        return combined[:2000]

    async def process_upload(
        self,
        filename: str,
        content: bytes,
        db: "AsyncSession | None" = None,
    ) -> dict:
        document_id, file_path = await self.save_file(filename, content)

        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "txt"
        type_map = {"pdf": "pdf", "txt": "txt", "docx": "docx", "pptx": "pptx"}
        file_type = type_map.get(ext, "txt")

        text = await self.extract_text(file_path, file_type)

        if not text.strip():
            raise DocumentProcessingError("File contains no extractable text")

        chunks = self.chunk_text(text)

        if db is not None:
            await self._persist_document(document_id, filename, file_path, text, chunks, db)
        else:
            self.documents[document_id] = {
                "filename": filename,
                "file_path": file_path,
                "chunks": chunks,
                "text": text,
            }

        return {
            "document_id": document_id,
            "filename": filename,
            "chunk_count": len(chunks),
        }

    async def _persist_document(
        self,
        document_id: str,
        filename: str,
        file_path: str,
        text: str,
        chunks: list[str],
        db: "AsyncSession",
    ) -> None:
        from app.db_models import Document

        db.add(Document(
            id=uuid.UUID(document_id),
            filename=filename,
            file_path=file_path,
            text=text,
            chunks=chunks,
            chunk_count=len(chunks),
        ))
        await db.flush()

    async def get_document(
        self,
        document_id: str,
        db: "AsyncSession | None" = None,
    ) -> dict | None:
        if db is not None:
            return await self._get_document_db(document_id, db)
        return self.documents.get(document_id)

    async def _get_document_db(
        self, document_id: str, db: "AsyncSession"
    ) -> dict | None:
        from sqlalchemy import select
        from app.db_models import Document

        result = await db.execute(
            select(Document).where(Document.id == uuid.UUID(document_id))
        )
        row = result.scalar_one_or_none()
        if row is None:
            return None
        return {
            "filename": row.filename,
            "file_path": row.file_path,
            "chunks": row.chunks,
            "text": row.text,
        }


document_service = DocumentService()
